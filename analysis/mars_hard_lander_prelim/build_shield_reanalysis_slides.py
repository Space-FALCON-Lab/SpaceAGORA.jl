from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

PPTX_DEPS = Path("/tmp/aero08_pptx_deps")
if PPTX_DEPS.exists():
    sys.path.insert(0, str(PPTX_DEPS))

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(
    "/Users/josephine/Research/Space-FALCON/_Proposals/[FY27]SURP/SHIELD_Reanalysis_Preliminary_Results"
)
PLOTS = ROOT / "plots"
TEMPLATE_ROOT = Path(
    "/Users/josephine/Research/Space-FALCON/_Command_Center/Templates/SpaceFALCON_UMich_Presentation_Template"
)
SITE_ASSETS = TEMPLATE_ROOT / "site_assets"
FALCON_LOGO = SITE_ASSETS / "space_falcon_logo.png"
UM_LOGO = SITE_ASSETS / "UM_Logo.png"
HERO_IMAGE = SITE_ASSETS / "FALCON_LAB_Home.png"

DECK_BASENAME = "SHIELD_Tuesday_Reanalysis_Results"
PPTX_PATH = ROOT / f"{DECK_BASENAME}.pptx"
OUTLINE_PATH = ROOT / f"{DECK_BASENAME}_outline.md"
QLPREVIEW_DIR = ROOT / f"{DECK_BASENAME}_qlpreview"

TITLE = RGBColor(17, 24, 39)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(75, 85, 99)
ACCENT = RGBColor(37, 99, 235)
ACCENT_FILL = RGBColor(239, 246, 255)
BORDER = RGBColor(191, 219, 254)
PANEL_FILL = RGBColor(248, 250, 252)
BG = RGBColor(255, 255, 255)
SITE_NAVY = RGBColor(11, 31, 58)
SITE_SKY = RGBColor(142, 198, 255)
UM_MAIZE = RGBColor(255, 203, 5)
DARK_BG = RGBColor(5, 7, 11)
SOFT_DARK = RGBColor(12, 17, 26)
OFF_WHITE = RGBColor(248, 250, 252)
WARN_FILL = RGBColor(254, 242, 242)
WARN_LINE = RGBColor(252, 165, 165)
WARN_TEXT = RGBColor(185, 28, 28)


def inches(x: float) -> int:
    return Inches(x)


def clean_outputs() -> None:
    for path in (PPTX_PATH, OUTLINE_PATH):
        if path.exists():
            path.unlink()
    if QLPREVIEW_DIR.exists():
        shutil.rmtree(QLPREVIEW_DIR)


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int,
    *,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = TEXT,
    font_name: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin_pt: float = 0,
):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(margin_pt)
    tf.margin_right = Pt(margin_pt)
    tf.margin_top = Pt(margin_pt)
    tf.margin_bottom = Pt(margin_pt)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    font.name = font_name
    return box


def add_bullets(slide, x, y, w, h, items: list[str], size: int = 23, color: RGBColor = TEXT):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            font = run.font
            font.size = Pt(size)
            font.color.rgb = color
            font.name = "Aptos"
    return box


def add_panel(slide, x, y, w, h, *, fill=PANEL_FILL, line=BORDER):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def add_image(slide, path: Path, x: float, y: float, *, w: float | None = None, h: float | None = None):
    if not path.exists():
        return
    kwargs = {}
    if w is not None:
        kwargs["width"] = inches(w)
    if h is not None:
        kwargs["height"] = inches(h)
    slide.shapes.add_picture(str(path), inches(x), inches(y), **kwargs)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float):
    from PIL import Image

    if not path.exists():
        return
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        final_w = w
        final_h = w / img_ratio
        final_x = x
        final_y = y + (h - final_h) / 2
    else:
        final_h = h
        final_w = h * img_ratio
        final_y = y
        final_x = x + (w - final_w) / 2
    slide.shapes.add_picture(str(path), inches(final_x), inches(final_y), width=inches(final_w), height=inches(final_h))


def add_brandbar(slide):
    add_panel(slide, 0, 0, 13.333, 0.64, fill=DARK_BG, line=DARK_BG)
    add_panel(slide, 0, 0.62, 13.333, 0.02, fill=UM_MAIZE, line=UM_MAIZE)
    add_image(slide, UM_LOGO, 0.5, 0.11, h=0.34)
    add_image(slide, FALCON_LOGO, 2.08, 0.02, h=0.58)
    add_textbox(
        slide,
        2.8,
        0.16,
        4.7,
        0.22,
        "Space-Flight Autonomous Leading CONcepts (Space-FALCON) Lab",
        11,
        color=OFF_WHITE,
        bold=True,
    )


def add_title(slide, title: str, subtitle: str):
    add_brandbar(slide)
    add_textbox(slide, 0.65, 0.92, 11.9, 0.48, title, 25, bold=True, color=SITE_NAVY)
    add_textbox(slide, 0.67, 1.36, 11.7, 0.32, subtitle, 11, color=MUTED)
    add_panel(slide, 0.67, 1.62, 0.95, 0.04, fill=UM_MAIZE, line=UM_MAIZE)
    add_panel(slide, 1.74, 1.62, 1.12, 0.03, fill=SITE_SKY, line=SITE_SKY)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def set_slide_notes(slide, notes_text: str):
    notes_slide = slide.notes_slide
    for placeholder in notes_slide.placeholders:
        if placeholder.placeholder_format.idx == 3:
            text_frame = placeholder.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = notes_text
            return


def add_stat_card(slide, x, y, w, h, heading: str, body: list[str], *, big: str | None = None, fill=ACCENT_FILL):
    add_panel(slide, x, y, w, h, fill=fill, line=BORDER if fill == ACCENT_FILL else fill)
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.36, 0.3, heading, 16, bold=True, color=ACCENT)
    body_y = y + 0.48
    if big is not None:
        add_textbox(slide, x + 0.18, body_y, w - 0.36, 0.7, big, 24, bold=True, color=TITLE)
        body_y += 0.72
    add_bullets(slide, x + 0.18, body_y, w - 0.36, h - (body_y - y) - 0.18, body, size=15)


def fmt_km(value: float) -> str:
    return f"{value:.0f}" if abs(value - round(value)) < 1e-9 else f"{value:.1f}"


SUMMARY = pd.read_csv(ROOT / "prelim_summary.csv")
AUTH = pd.read_csv(ROOT / "authority_summary.csv").iloc[0]
TARGET = pd.read_csv(ROOT / "target_guidance_summary.csv").iloc[0]
MC = pd.read_csv(ROOT / "monte_carlo_summary.csv")
CROSS = pd.read_csv(ROOT / "crossrange_sensitivity.csv")
ALPHA = pd.read_csv(ROOT / "alpha_body_sensitivity.csv")
LOCAL = pd.read_csv(ROOT / "local_control_effectiveness.csv")

BODY = SUMMARY[SUMMARY["policy"] == "body_only"].iloc[0]
DEPLOYED = SUMMARY[SUMMARY["policy"] == "fixed_deployed"].iloc[0]
BANG = SUMMARY[SUMMARY["policy"] == "bang_bang"].sort_values("h_switch_km")
BEST_CROSS = CROSS.sort_values("abs_impact_crossrange_km", ascending=False).iloc[0]
GUIDED = MC[MC["policy"] == "guided_targeted_optimistic"].iloc[0]
CROSS_GUIDED = MC[MC["policy"] == "guided_targeted_crossrange_bilateral"].iloc[0]
CONE_HALF_ANGLE_DEG = 50
SWITCH_BAND_MIN_KM = float(BANG["h_switch_km"].min())
SWITCH_BAND_MAX_KM = float(BANG["h_switch_km"].max())
PEAK_EFFECTIVENESS_KM = float(LOCAL.loc[LOCAL["sensitivity_abs_km_per_km"].idxmax(), "h_switch_km"])


def slide1(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_image_fit(slide, HERO_IMAGE, 0.0, 0.0, 13.333, 7.5)
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inches(0), inches(0), inches(13.333), inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DARK_BG
    overlay.fill.transparency = 0.18
    overlay.line.fill.background()
    add_panel(slide, 0, 0, 13.333, 0.18, fill=UM_MAIZE, line=UM_MAIZE)
    add_image(slide, UM_LOGO, 0.78, 0.38, h=0.42)
    add_image(slide, FALCON_LOGO, 2.12, 0.18, h=0.9)
    add_textbox(slide, 0.92, 1.24, 5.4, 0.22, "SPACE-FALCON LAB", 12, bold=True, color=SITE_SKY)
    add_textbox(slide, 0.9, 1.58, 6.1, 0.9, "Published-Size SHIELD Reanalysis", 30, bold=True, color=OFF_WHITE, font_name="Aptos Display")
    add_textbox(slide, 0.92, 2.42, 6.0, 0.48, "First-cut guidance study using subsonic deployment timing on a published-size SHIELD surrogate.", 18, color=RGBColor(226, 232, 240))
    add_panel(slide, 0.92, 3.1, 0.95, 0.08, fill=UM_MAIZE, line=UM_MAIZE)
    add_panel(slide, 1.98, 3.12, 1.45, 0.05, fill=SITE_SKY, line=SITE_SKY)
    add_bullets(
        slide,
        0.92,
        3.55,
        5.45,
        1.7,
        [
            "Published SHIELD basis: 1.8 m stowed diameter, > 2 m deployed drag surface, 120 kg entry mass.",
            f"Current model: {CONE_HALF_ANGLE_DEG}° sphere-cone surrogate + drag-skirt-equivalent deployed area derived from the 2.2 m / 0.75 m skirt surrogate.",
            "Goal: quantify downrange authority from deployment timing, not claim flight-truth 6DOF guidance.",
        ],
        size=17,
        color=OFF_WHITE,
    )
    add_panel(slide, 0.92, 5.65, 5.3, 1.02, fill=SOFT_DARK, line=SOFT_DARK)
    add_textbox(slide, 1.18, 5.88, 4.8, 0.2, "Representative result", 14, bold=True, color=SITE_SKY)
    add_textbox(slide, 1.18, 6.13, 4.8, 0.36, f"{AUTH['downrange_authority_km']:.1f} km downrange authority", 24, bold=True, color=UM_MAIZE, font_name="Aptos Display")
    add_textbox(
        slide,
        1.18,
        6.48,
        4.8,
        0.2,
        f"120 kg entry mass   |   β_high = {AUTH['target_beta_high_kg_m2']:.1f} kg/m²   |   β_low = {AUTH['target_beta_low_kg_m2']:.1f} kg/m²",
        11,
        color=RGBColor(191, 219, 254),
    )
    add_panel(slide, 7.05, 1.45, 5.45, 5.2, fill=RGBColor(255, 255, 255), line=RGBColor(214, 229, 248))
    add_textbox(slide, 7.38, 1.78, 4.9, 0.34, "Published SHIELD values used in the rerun", 18, bold=True, color=SITE_NAVY)
    add_bullets(
        slide,
        7.4,
        2.2,
        4.85,
        3.6,
        [
            "Payload up to 6 kg",
            "Stowed diameter 1.8 m",
            "Drag skirt height 0.75 m",
            "Entry mass 120 kg, landed mass 65 kg (Table 1 basis)",
            "Deployed ballistic coefficient target < 10 kg/m²",
            "Published terminal velocity target ≤ 70 m/s",
        ],
        size=16,
    )
    add_panel(slide, 7.35, 5.92, 4.9, 0.52, fill=ACCENT_FILL, line=BORDER)
    add_textbox(slide, 7.58, 6.09, 4.45, 0.18, "This rerun is for the right vehicle class; the old β = 150 surrogate was not.", 12, color=ACCENT, bold=True)


def slide2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Subsonic Deployment Timing Produces Meaningful Downrange Authority",
        f"For the published-size SHIELD surrogate, switch timing across the {fmt_km(SWITCH_BAND_MAX_KM)}–{fmt_km(SWITCH_BAND_MIN_KM)} km subsonic deployment band opens an {AUTH['downrange_authority_km']:.0f} km reachable downrange corridor.",
    )
    add_image_fit(slide, PLOTS / "downrange_authority_vs_switch_altitude.png", 0.35, 1.86, 8.0, 5.2)
    add_panel(slide, 8.65, 1.95, 3.75, 3.12)
    add_bullets(
        slide,
        8.88,
        2.16,
        3.3,
        2.85,
        [
            f"Body only reaches {BODY['impact_downrange_km']:.1f} km; always-deployed reaches {DEPLOYED['impact_downrange_km']:.1f} km.",
            f"Reachable corridor = {AUTH['downrange_authority_km']:.1f} km.",
            f"Targeted solve converges at h_switch = {TARGET['h_switch_km']:.2f} km with {TARGET['range_error_km']:.3f} km error.",
            "This is passive skirt-deployment timing, not rim-flap actuation or hypersonic steering.",
        ],
        size=15,
    )
    add_stat_card(
        slide,
        8.65,
        5.35,
        3.75,
        1.35,
        "Answer to the scoping question",
        [f"This first-cut SHIELD-class rerun lands in the “~{AUTH['downrange_authority_km']:.0f} km” authority regime."],
    )


def slide3(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Authority Is Generated In A Narrow Subsonic Corridor",
        "The effect is localized near the deployment band rather than spread uniformly across entry.",
    )
    add_image_fit(slide, PLOTS / "local_control_effectiveness.png", 0.35, 1.85, 8.2, 4.2)
    add_image_fit(slide, PLOTS / "drag_accel_vs_altitude_clean.png", 8.82, 1.98, 3.63, 2.62)
    add_panel(slide, 8.82, 4.88, 3.63, 1.72)
    add_bullets(
        slide,
        9.04,
        5.06,
        3.2,
        1.45,
        [
            f"Sensitivity peaks near h_switch ≈ {fmt_km(PEAK_EFFECTIVENESS_KM)} km, close to the upper end of the subsonic band.",
            "That same band is where deployed drag begins to reshape the deceleration history.",
            "This supports a simple switch-altitude guidance law.",
        ],
        size=14,
    )
    set_slide_notes(
        slide,
        "The left plot shows signed range sensitivity to switch altitude. The right support plot compares only three trajectories: body only, always deployed to impact, and the targeted switch solution. The point is that the authority comes from reshaping the drag history through the subsonic deployment band, not from a diffuse whole-entry effect.",
    )


def slide4(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Deployment Timing Trades Range Against Impact Conditions",
        "Earlier deployment buys range reduction, but it also changes impact velocity and peak aero load.",
    )
    add_image_fit(slide, PLOTS / "terminal_velocity_and_peak_g_vs_switch_altitude.png", 0.35, 1.86, 8.15, 5.0)
    add_panel(slide, 8.75, 1.95, 3.7, 3.2)
    add_bullets(
        slide,
        8.98,
        2.15,
        3.2,
        2.9,
        [
            f"Earliest allowed deployment lowers impact velocity to about {DEPLOYED['impact_velocity_mps']:.0f} m/s.",
            f"Late / no deployment rises to about {BODY['impact_velocity_mps']:.0f} m/s.",
            f"Peak aero load reaches about {BANG['peak_total_decel_earth_g'].max():.0f} g for the earliest deployment cases.",
            "Published SHIELD targets ≤ 70 m/s at impact / terminal speed, so this surrogate is still conservative on deployed drag performance.",
        ],
        size=14,
    )
    add_panel(slide, 8.75, 5.45, 3.7, 1.1, fill=ACCENT_FILL, line=BORDER)
    add_textbox(slide, 8.98, 5.67, 3.15, 0.18, "Interpretation", 14, bold=True, color=ACCENT)
    add_textbox(slide, 8.98, 5.93, 3.15, 0.42, "This slide is still symmetric drag modulation only, so it is a downrange result, not a crossrange result.", 12, color=TEXT)
    set_slide_notes(
        slide,
        "This is the main trade slide for the SHIELD rerun. The new analysis is no longer about β-ratio sweeps; it is about what deployment timing does to reachable range and impact conditions for the published-size case. The caveat is important: the surrogate still misses the published ≤70 m/s impact-speed target, which tells us the drag-skirt representation is still conservative.",
    )


def slide5(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Crossrange Still Requires Asymmetry",
        "Preliminary crossrange appears only when the left and right deployed surfaces no longer match, and it remains strongly coupled to downrange.",
    )
    add_image_fit(slide, PLOTS / "crossrange_sensitivity.png", 0.35, 1.82, 7.8, 5.2)
    add_panel(slide, 8.35, 1.92, 4.25, 2.58)
    add_bullets(
        slide,
        8.58,
        2.12,
        3.8,
        2.35,
        [
            f"Best deterministic crossrange in this sweep: {BEST_CROSS['abs_impact_crossrange_km']:.1f} km.",
            f"That best case also shifts downrange by {BEST_CROSS['downrange_shift_vs_symmetric_km']:.1f} km.",
            f"Crossrange-aware 95% ellipse: {CROSS_GUIDED['ellipse_major_axis_km']:.1f} km × {CROSS_GUIDED['ellipse_minor_axis_km']:.1f} km.",
            "With 2 panels, crossrange is real but not cleanly decoupled yet.",
        ],
        size=14,
    )
    add_image_fit(slide, PLOTS / "landing_ellipse_centered.png", 8.45, 4.75, 4.05, 1.95)
    add_textbox(
        slide,
        8.55,
        6.78,
        3.95,
        0.24,
        "This is still a preliminary differential-panel extension, not a mature SHIELD guidance claim.",
        11,
        color=MUTED,
    )
    set_slide_notes(
        slide,
        "The crossrange message is intentionally cautious. The SHIELD-class rerun still shows that downrange authority is the strongest result. Crossrange exists only when left and right deployment differ, and with a 2-panel architecture it remains strongly coupled to range.",
    )


def slide6(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Backup: Modeling Setup And Initial Conditions",
        "What the published-size SHIELD rerun includes, what it excludes, and how the deployment band is treated.",
    )
    add_stat_card(
        slide,
        0.55,
        1.92,
        3.75,
        2.05,
        "Nominal entry state",
        [
            "h₀ = 125 km",
            "V₀ = 5.5 km/s",
            "γ₀ = -12°",
            "Latitude = 0°",
        ],
    )
    add_stat_card(
        slide,
        0.55,
        4.15,
        3.75,
        2.45,
        "Published SHIELD basis used here",
        [
            "1.8 m stowed diameter",
            "> 2 m deployed drag surface",
            "0.75 m drag-skirt height",
            "120 kg entry mass, 65 kg landed mass",
        ],
    )
    add_panel(slide, 4.6, 1.92, 3.7, 4.7)
    add_textbox(slide, 4.84, 2.14, 3.2, 0.34, "Included in this first cut", 18, bold=True, color=ACCENT)
    add_bullets(
        slide,
        4.82,
        2.5,
        3.2,
        3.95,
        [
            "MarsGRAM-backed 3DOF / point-mass entry model",
            "Fixed 120 kg entry mass with stowed β calibrated from the surrogate body",
            f"Subsonic deployment timing sweep over {fmt_km(SWITCH_BAND_MAX_KM)}–{fmt_km(SWITCH_BAND_MIN_KM)} km",
            "Main deterministic authority plots use zero winds",
            "3D Monte Carlo footprint uses MarsGRAM winds",
        ],
        size=14,
    )
    add_panel(slide, 8.65, 1.92, 3.75, 4.7, fill=WARN_FILL, line=WARN_LINE)
    add_textbox(slide, 8.9, 2.14, 3.2, 0.34, "Not included yet", 18, bold=True, color=WARN_TEXT)
    add_bullets(
        slide,
        8.88,
        2.5,
        3.2,
        3.95,
        [
            "Exact SHIELD drag-skirt geometry or deployment kinematics",
            "Full 6DOF attitude dynamics",
            "Trim, stability, hinge-load, or structural closure",
            "CFD / DSMC or transition-flow truth modeling",
            "Precision landing guidance",
        ],
        size=14,
    )


def slide7(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Backup: Published SHIELD Basis vs Current Surrogate Mapping",
        "What came from Barba 2021, and how it was translated into the current first-cut rerun.",
    )
    add_panel(slide, 0.55, 1.88, 5.8, 4.95)
    add_textbox(slide, 0.82, 2.08, 5.2, 0.3, "Published SHIELD numbers", 20, bold=True, color=ACCENT)
    add_bullets(
        slide,
        0.82,
        2.44,
        5.05,
        4.15,
        [
            "Payload up to 6 kg",
            "Stowed diameter 1.8 m",
            "Deployed drag surface > 2 m diameter",
            "Drag-skirt height 0.75 m",
            "Table 1 basis used here: 120 kg entry mass, 65 kg landed mass",
            "Published deployed ballistic coefficient target < 10 kg/m²",
            "Published terminal velocity target ≤ 70 m/s",
            "Published 99% ellipse 40 km × 2 km",
        ],
        size=15,
    )
    add_panel(slide, 6.82, 1.88, 5.95, 4.95)
    add_textbox(slide, 7.09, 2.08, 5.35, 0.3, "Current surrogate interpretation", 20, bold=True, color=ACCENT)
    add_bullets(
        slide,
        7.08,
        2.44,
        5.2,
        4.15,
        [
            f"{CONE_HALF_ANGLE_DEG}° sphere-cone surrogate with base radius 0.9 m and nose radius 0.288 m",
            f"Fixed 120 kg entry mass gives stowed β_high = {AUTH['target_beta_high_kg_m2']:.1f} kg/m²",
            f"Architecture-derived deployed β_low = {AUTH['target_beta_low_kg_m2']:.1f} kg/m²",
            f"Equivalent deployed area = {BODY['panel_area_total_m2']:.2f} m² total ({BODY['panel_area_each_m2']:.2f} m² per panel), derived from the 2.2 m / 0.75 m skirt surrogate",
            f"Illustrative panel-system mass at 5 kg/m² = {BODY['panel_area_total_m2'] * 5.0:.1f} kg",
            f"Subsonic deployment timing sweep from {fmt_km(SWITCH_BAND_MAX_KM)} km to {fmt_km(SWITCH_BAND_MIN_KM)} km",
            "Directionally useful for guidance authority, but not yet a true drag-skirt fidelity model",
        ],
        size=15,
    )
    add_panel(slide, 0.78, 6.15, 11.95, 0.58, fill=ACCENT_FILL, line=BORDER)
    add_textbox(slide, 1.02, 6.33, 11.4, 0.18, "Most important change from the old deck: this rerun is now tied to the published SHIELD vehicle class rather than the discarded β = 150 surrogate.", 13, color=ACCENT, bold=True)


def slide8(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Backup: Open Risks And Next Gating Questions",
        "The main remaining questions are fidelity and engineering closure questions, not whether the concept produces any authority at all.",
    )
    add_panel(slide, 0.58, 1.88, 5.95, 4.98)
    add_textbox(slide, 0.85, 2.1, 5.35, 0.3, "What still needs to close", 20, bold=True, color=ACCENT)
    add_bullets(
        slide,
        0.84,
        2.46,
        5.3,
        4.1,
        [
            "Drag-skirt fidelity: the surrogate still lands above the published ≤70 m/s impact-speed target in the fully deployed state.",
            "Trim and stability: deployment changes moments and center-of-pressure location, not just drag.",
            "Thermal and structural survivability: skirt deployment, loads, and impact compatibility still need closure.",
            "Mass closure: the equivalent panel area is large enough that hardware packaging matters immediately.",
        ],
        size=15,
    )
    add_panel(slide, 6.78, 1.88, 5.95, 4.98)
    add_textbox(slide, 7.05, 2.1, 5.35, 0.3, "What the next iteration should answer", 20, bold=True, color=ACCENT)
    add_bullets(
        slide,
        7.04,
        2.46,
        5.3,
        4.1,
        [
            "Can a more faithful deployed drag-skirt model recover the published low-β / low-impact-speed behavior while preserving the authority trend?",
            "How much authority survives once the true SHIELD geometry and deployment mechanics are modeled more directly?",
            "Is the best proposal framing “SHIELD with steering” or “SHIELD-class follow-on with active drag modulation”?",
            "What minimum asymmetry or extra surface count would be needed before crossrange becomes mission-relevant?",
        ],
        size=15,
    )
    add_panel(slide, 0.72, 6.15, 12.0, 0.62, fill=ACCENT_FILL, line=BORDER)
    add_textbox(
        slide,
        0.96,
        6.34,
        11.4,
        0.18,
        f"Current proposal-safe claim: published-size SHIELD-class deployment timing appears to buy ~{AUTH['downrange_authority_km']:.0f} km of downrange authority in this first-cut surrogate model; crossrange remains preliminary.",
        13,
        bold=True,
        color=ACCENT,
    )


def build_outline():
    OUTLINE_PATH.write_text(
        """# SHIELD Tuesday Reanalysis Outline

1. Published-size SHIELD reanalysis cover and scope reset
2. Downrange authority from subsonic deployment timing
3. Where the authority is generated
4. Impact-condition trade vs switch altitude
5. Crossrange as an asymmetric extension
6. Backup: modeling setup and initial conditions
7. Backup: published SHIELD basis vs surrogate mapping
8. Backup: open risks and next gating questions
"""
    )


def export_quicklook_preview() -> None:
    subprocess.run(
        ["qlmanage", "-p", "-o", str(QLPREVIEW_DIR.parent), str(PPTX_PATH)],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    generated = QLPREVIEW_DIR.parent / f"{PPTX_PATH.name}.qlpreview"
    if generated.exists() and generated != QLPREVIEW_DIR:
        if QLPREVIEW_DIR.exists():
            shutil.rmtree(QLPREVIEW_DIR)
        generated.rename(QLPREVIEW_DIR)


def main():
    clean_outputs()
    prs = Presentation()
    prs.slide_width = inches(13.333)
    prs.slide_height = inches(7.5)
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    prs.save(str(PPTX_PATH))
    build_outline()
    export_quicklook_preview()
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
