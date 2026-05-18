from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

PPTX_DEPS = Path("/tmp/aero08_pptx_deps")
if PPTX_DEPS.exists():
    sys.path.insert(0, str(PPTX_DEPS))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROPOSAL_ROOT = Path(
    "/Users/josephine/Research/Space-FALCON/_Proposals/[FY27]SURP/AERO-08_Preliminary_Results"
)
PLOTS_ROOT = PROPOSAL_ROOT / "plots"
TEMPLATE_ROOT = Path(
    "/Users/josephine/Research/Space-FALCON/_Command_Center/Templates/SpaceFALCON_UMich_Presentation_Template"
)
SITE_ASSETS = TEMPLATE_ROOT / "site_assets"
FALCON_LOGO = SITE_ASSETS / "space_falcon_logo.png"
UM_LOGO = SITE_ASSETS / "UM_Logo.png"
HERO_IMAGE = SITE_ASSETS / "FALCON_LAB_Home.png"
DECK_BASENAME = "AERO-08_Tuesday_Preliminary_Results"
PPTX_PATH = PROPOSAL_ROOT / f"{DECK_BASENAME}.pptx"
PDF_PATH = PROPOSAL_ROOT / f"{DECK_BASENAME}.pdf"
IMAGES_DIR = PROPOSAL_ROOT / f"{DECK_BASENAME}_images"
QLPREVIEW_DIR = PROPOSAL_ROOT / f"{DECK_BASENAME}_qlpreview"
OUTLINE_PATH = PROPOSAL_ROOT / f"{DECK_BASENAME}_outline.md"

TITLE = RGBColor(17, 24, 39)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(75, 85, 99)
ACCENT = RGBColor(37, 99, 235)
ACCENT_FILL = RGBColor(239, 246, 255)
BORDER = RGBColor(191, 219, 254)
PANEL_FILL = RGBColor(248, 250, 252)
BG = RGBColor(255, 255, 255)
SITE_NAVY = RGBColor(11, 31, 58)
SITE_SKY = RGBColor(142, 198, 255)
UM_MAIZE = RGBColor(255, 203, 5)
DARK_BG = RGBColor(5, 7, 11)
SOFT_DARK = RGBColor(12, 17, 26)
OFF_WHITE = RGBColor(248, 250, 252)


def inches(x: float) -> int:
    return Inches(x)


def clean_outputs() -> None:
    for path in (
        PPTX_PATH,
        PDF_PATH,
        OUTLINE_PATH,
        PROPOSAL_ROOT / f"{DECK_BASENAME}.key",
        PROPOSAL_ROOT / f"~${DECK_BASENAME}.pptx",
    ):
        if path.exists():
            path.unlink()
    for path in (IMAGES_DIR, QLPREVIEW_DIR):
        if path.exists():
            shutil.rmtree(path)


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int,
    *,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = TEXT,
    font_name: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin_pt: float = 0,
):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(margin_pt)
    tf.margin_right = Pt(margin_pt)
    tf.margin_top = Pt(margin_pt)
    tf.margin_bottom = Pt(margin_pt)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    font.name = font_name
    return box


def add_bullets(slide, x, y, w, h, items: list[str], size: int = 23):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            font = run.font
            font.size = Pt(size)
            font.color.rgb = TEXT
            font.name = "Aptos"
    return box


def add_panel(slide, x, y, w, h, *, fill=PANEL_FILL, line=BORDER):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.4)
    return shape


def add_stat_card(slide, x, y, w, h, heading: str, body: list[str], *, big: str | None = None):
    add_panel(slide, x, y, w, h, fill=ACCENT_FILL, line=BORDER)
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.36, 0.35, heading, 16, bold=True, color=ACCENT)
    body_y = y + 0.48
    if big is not None:
        add_textbox(slide, x + 0.18, body_y, w - 0.36, 0.7, big, 26, bold=True, color=TITLE)
        body_y += 0.78
    add_bullets(slide, x + 0.18, body_y, w - 0.36, h - (body_y - y) - 0.18, body, size=16)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        final_w = w
        final_h = w / img_ratio
        final_x = x
        final_y = y + (h - final_h) / 2
    else:
        final_h = h
        final_w = h * img_ratio
        final_y = y
        final_x = x + (w - final_w) / 2
    slide.shapes.add_picture(str(path), inches(final_x), inches(final_y), width=inches(final_w), height=inches(final_h))


def add_image(slide, path: Path, x: float, y: float, *, w: float | None = None, h: float | None = None):
    if not path.exists():
        return
    kwargs = {}
    if w is not None:
        kwargs["width"] = inches(w)
    if h is not None:
        kwargs["height"] = inches(h)
    slide.shapes.add_picture(str(path), inches(x), inches(y), **kwargs)


def add_brandbar(slide):
    add_panel(slide, 0, 0, 13.333, 0.64, fill=DARK_BG, line=DARK_BG)
    add_panel(slide, 0, 0.62, 13.333, 0.02, fill=UM_MAIZE, line=UM_MAIZE)
    add_image(slide, UM_LOGO, 0.5, 0.11, h=0.34)
    add_image(slide, FALCON_LOGO, 2.08, 0.02, h=0.58)
    add_textbox(
        slide,
        2.8,
        0.16,
        4.7,
        0.22,
        "Space-Flight Autonomous Leading CONcepts (Space-FALCON) Lab",
        11,
        color=OFF_WHITE,
        bold=True,
    )


def add_title(slide, title: str, subtitle: str):
    add_brandbar(slide)
    add_textbox(slide, 0.65, 0.92, 11.8, 0.48, title, 25, bold=True, color=SITE_NAVY)
    add_textbox(slide, 0.67, 1.36, 11.6, 0.28, subtitle, 11, color=MUTED)
    add_panel(slide, 0.67, 1.62, 0.95, 0.04, fill=UM_MAIZE, line=UM_MAIZE)
    add_panel(slide, 1.74, 1.62, 1.12, 0.03, fill=SITE_SKY, line=SITE_SKY)


def set_slide_notes(slide, notes_text: str):
    notes_slide = slide.notes_slide
    for placeholder in notes_slide.placeholders:
        if placeholder.placeholder_format.idx == 3:
            text_frame = placeholder.text_frame
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = notes_text
            return


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def slide1(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_image_fit(slide, HERO_IMAGE, 0.0, 0.0, 13.333, 7.5)
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, inches(0), inches(0), inches(13.333), inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DARK_BG
    overlay.fill.transparency = 0.18
    overlay.line.fill.background()
    add_panel(slide, 0, 0, 13.333, 0.18, fill=UM_MAIZE, line=UM_MAIZE)
    add_image(slide, UM_LOGO, 0.78, 0.38, h=0.42)
    add_image(slide, FALCON_LOGO, 2.12, 0.18, h=0.9)
    add_textbox(slide, 0.92, 1.26, 5.2, 0.22, "SPACE-FALCON LAB", 12, bold=True, color=SITE_SKY)
    add_textbox(slide, 0.9, 1.62, 6.1, 0.7, "Low-Cost Guided Mars Surface Access", 29, bold=True, color=OFF_WHITE, font_name="Aptos Display")
    add_textbox(slide, 0.92, 2.48, 5.8, 0.46, "Deployable aerodynamic area can turn a SHIELD-like hard-lander into a coarse-targeting Mars entry system.", 18, color=RGBColor(226, 232, 240))
    add_panel(slide, 0.92, 3.18, 0.95, 0.08, fill=UM_MAIZE, line=UM_MAIZE)
    add_panel(slide, 1.98, 3.2, 1.45, 0.05, fill=SITE_SKY, line=SITE_SKY)
    add_bullets(
        slide,
        0.92,
        3.7,
        5.35,
        1.75,
        [
            "Passive hard-lander body in the stowed state",
            "Discrete deploy / jettison drag control during entry",
            "Goal: footprint reduction, not precision landing",
        ],
        size=18,
    )
    add_panel(slide, 0.92, 5.72, 4.95, 0.95, fill=SOFT_DARK, line=SOFT_DARK)
    add_textbox(slide, 1.18, 5.96, 4.45, 0.24, "Representative result", 14, bold=True, color=SITE_SKY)
    add_textbox(slide, 1.18, 6.25, 4.45, 0.34, "103.9 km downrange authority", 24, bold=True, color=UM_MAIZE, font_name="Aptos Display")
    add_textbox(slide, 1.18, 6.58, 4.45, 0.24, "β_high = 150 kg/m²   |   β_ratio = 4   |   h_jettison = 20 km", 11, color=RGBColor(191, 219, 254))

    add_panel(slide, 6.85, 1.32, 5.95, 4.35, fill=RGBColor(255, 255, 255), line=RGBColor(214, 229, 248))
    add_image_fit(slide, PLOTS_ROOT / "body_schematic.png", 7.02, 1.48, 5.61, 4.01)
    add_textbox(
        slide,
        7.35,
        5.92,
        5.1,
        0.2,
        "Dark website-style cover with the concept figure and live lab branding.",
        10,
        italic=True,
        color=RGBColor(191, 219, 254),
        align=PP_ALIGN.RIGHT,
    )


def slide2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Drag Modulation Produces Meaningful Downrange Authority",
        "Authority grows with β-ratio and collapses if the surfaces must be jettisoned too early.",
    )
    add_image_fit(slide, PLOTS_ROOT / "downrange_authority_heatmap.png", 0.42, 1.86, 7.65, 5.05)
    add_panel(slide, 8.35, 1.95, 4.25, 3.05)
    add_bullets(
        slide,
        8.6,
        2.18,
        3.75,
        2.7,
        [
            "Nominal case: 103.9 km authority.",
            "β ratio is the main control-authority lever.",
            "Keeping the surfaces available below about 20 km matters.",
            "This is the direct answer to the “5 km or 50 km?” question.",
        ],
        size=19,
    )
    add_stat_card(
        slide,
        8.35,
        5.2,
        4.25,
        1.7,
        "Nominal case",
        ["β high = 150 kg/m²", "β ratio = 4", "h_jettison = 20 km"],
    )


def slide3(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Authority Is Generated In A Specific Entry Corridor",
        "Switch timing matters most in a mid-entry corridor where the atmosphere becomes useful.",
    )
    add_image_fit(slide, PLOTS_ROOT / "local_control_effectiveness.png", 0.35, 1.85, 8.35, 4.2)
    add_image_fit(slide, PLOTS_ROOT / "drag_accel_vs_altitude.png", 8.95, 1.98, 3.45, 2.2)
    add_panel(slide, 8.95, 4.28, 3.45, 2.3)
    add_bullets(
        slide,
        9.15,
        4.47,
        3.05,
        2.2,
        [
            "Control is physically localized, not numerically smeared across the whole entry.",
            "The same corridor is where drag acceleration rises sharply.",
            "This supports a discrete switch-altitude guidance logic.",
        ],
        size=15,
    )
    set_slide_notes(
        slide,
        "Figure label guide for the small drag-acceleration support plot: "
        "'Body only' means the passive entry body with no deployed surfaces. "
        "'Always deployed until 20 km' means the drag panels stay deployed from entry until jettison at 20 km altitude. "
        "'Deploy below X km, jettison at 20 km' means the panels stay stowed above the named switch altitude, deploy once the vehicle descends below that altitude, and are then jettisoned at 20 km. "
        "These labels describe the discrete aerodynamic control schedule rather than different vehicle shapes.",
    )


def slide4(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Control Authority Comes With A Design Trade",
        "Slides 2-4 use symmetric deploy / jettison drag modulation only, so these results quantify downrange authority rather than crossrange steering.",
    )
    add_image_fit(slide, PLOTS_ROOT / "terminal_velocity_and_peak_g_vs_beta_ratio.png", 0.35, 1.85, 8.4, 4.3)
    add_panel(slide, 0.55, 6.34, 8.1, 0.48, fill=ACCENT_FILL, line=BORDER)
    add_textbox(
        slide,
        0.78,
        6.47,
        7.6,
        0.18,
        "Interpretation: more symmetric drag area changes how far the vehicle flies downrange, but it does not create independent lateral control.",
        12,
        color=ACCENT,
        bold=True,
    )
    add_panel(slide, 8.95, 1.92, 3.45, 3.45)
    add_bullets(
        slide,
        9.15,
        2.14,
        3.05,
        3.25,
        [
            "β ratio = β_high / β_low.",
            "Nominal β ratio = 4 means 0.52 m² total panel area, or 0.26 m² per panel.",
            "At 5 kg/m² system areal density, that is about 2.6 kg of added panel system mass.",
            "β ratio ≈ 8 is stronger, but looks aggressive in area and mechanism burden.",
        ],
        size=14,
    )
    add_stat_card(
        slide,
        8.95,
        5.6,
        3.45,
        1.0,
        "Current recommendation",
        ["Lead with β ratio ≈ 4 as the current sweet spot."],
    )
    set_slide_notes(
        slide,
        "This slide is still part of the symmetric drag-modulation story. "
        "The deployed surfaces are matched left-right, so the only intentional control authority here is through total drag area and therefore downrange shaping. "
        "Crossrange is introduced only on the next slide, when the left and right panels are allowed to differ.",
    )


def slide5(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Crossrange Requires Asymmetric Deployment",
        "Preliminary crossrange appears only when the left and right panels no longer match; with 2 panels it remains strongly coupled to downrange.",
    )
    add_image_fit(slide, PLOTS_ROOT / "crossrange_sensitivity.png", 0.35, 1.82, 7.8, 5.2)
    add_panel(slide, 8.35, 1.92, 4.25, 2.45)
    add_bullets(
        slide,
        8.6,
        2.12,
        3.8,
        2.25,
        [
            "This is a separate extension beyond the symmetric drag-only cases on slides 2-4.",
            "Best deterministic crossrange offset in this sweep: 18.7 km.",
            "With 2 panels, crossrange comes with a large downrange perturbation.",
            "This is a promising extension, not the main demonstrated claim yet.",
        ],
        size=15,
    )
    add_image_fit(slide, PLOTS_ROOT / "landing_ellipse_centered.png", 8.5, 4.62, 4.0, 2.12)
    add_textbox(
        slide,
        8.55,
        6.8,
        3.95,
        0.28,
        "Future work: more panels or a different architecture may decouple crossrange from total drag.",
        11,
        color=MUTED,
    )
    set_slide_notes(
        slide,
        "Crossrange is not part of the symmetric drag-only results on slides 2 through 4. "
        "It appears only when the left and right panels are intentionally mismatched in deployment, cant, or deflection, which creates a net side force. "
        "The point of this slide is to show that crossrange is possible, but still coupled to downrange with the current 2-panel architecture.",
    )


def slide6(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Backup: Modeling Setup And Initial Conditions",
        "What the preliminary analysis includes, what it excludes, and the nominal entry state.",
    )
    add_stat_card(
        slide,
        0.55,
        1.92,
        3.75,
        1.95,
        "Nominal entry state",
        [
            "h₀ = 125 km",
            "V₀ = 5.5 km/s",
            "γ₀ = -12°",
            "Latitude = 0°",
        ],
    )
    add_stat_card(
        slide,
        0.55,
        4.05,
        3.75,
        2.6,
        "Dynamics and atmosphere",
        [
            "Mars spherical gravity",
            "Planar 3DOF point-mass entry model",
            "MarsGRAM density and temperature for the main sweeps",
            "Discrete body-only / deployed / jettison states",
        ],
    )
    add_panel(slide, 4.6, 1.92, 3.7, 4.2)
    add_textbox(slide, 4.84, 2.14, 3.2, 0.34, "Included in this first cut", 18, bold=True, color=ACCENT)
    add_bullets(
        slide,
        4.82,
        2.5,
        3.2,
        3.35,
        [
            "Deterministic downrange sweeps over β ratio, jettison altitude, and switch altitude",
            "Main authority plots use zero winds to isolate aerodynamic control",
            "Target-range solve over h_switch",
            "3D Monte Carlo footprint uses MarsGRAM winds",
            "Preliminary asymmetric-panel crossrange sensitivity",
        ],
        size=15,
    )
    add_panel(slide, 8.65, 1.92, 3.75, 4.2, fill=RGBColor(254, 242, 242), line=RGBColor(252, 165, 165))
    add_textbox(slide, 8.9, 2.14, 3.2, 0.34, "Not included yet", 18, bold=True, color=RGBColor(185, 28, 28))
    add_bullets(
        slide,
        8.88,
        2.5,
        3.2,
        3.35,
        [
            "Full 6DOF attitude dynamics",
            "DSMC / CFD or transition-flow truth modeling",
            "Trim, stability, or hinge-load analysis",
            "Precision landing guidance",
        ],
        size=15,
    )
    set_slide_notes(
        slide,
        "MarsGRAM does include winds in this workflow. "
        "The deterministic downrange authority sweeps are intentionally run with zero winds so the main plots isolate aerodynamic control authority rather than weather dispersion. "
        "The wind-aware path is used in the 3D Monte Carlo landing-ellipse calculation.",
    )


def slide7(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Backup: Vehicle And Aerodynamic Construction",
        "SHIELD-informed surrogate body and deployable-surface assumptions used in the current study.",
    )
    add_image_fit(slide, PLOTS_ROOT / "body_schematic.png", 0.48, 1.82, 12.35, 5.35)
    set_slide_notes(
        slide,
        "This backup figure is meant to do the geometry explanation visually. "
        "The left column shows the passive body-only state and the symmetric deployed drag state. "
        "The right column shows the top-view panel symmetry, the side-view analysis geometry, and the aerodynamic assumptions carried in the current model.",
    )


def slide8(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Backup: What β Ratio Means In Hardware",
        "Same body mass class, different deployed panel area and added panel-system mass.",
    )
    add_textbox(
        slide,
        0.75,
        1.92,
        7.5,
        0.52,
        "Nominal reference case: β_high = 150 kg/m², entry mass ≈ 48.46 kg, 2 symmetric deployed panels.",
        17,
        color=TEXT,
    )
    add_textbox(
        slide,
        0.75,
        2.32,
        8.2,
        0.44,
        "Interpretation: β ratio = β_high / β_low, so larger ratio means more deployed drag area and stronger control authority.",
        15,
        color=MUTED,
    )

    x0 = 0.78
    widths = [1.45, 2.05, 2.05, 2.0, 2.0, 1.65]
    headers = ["β ratio", "Total panel area", "Area per panel", "Added panel mass*", "Authority", "Comment"]
    y_head = 2.95
    row_h = 0.78
    x = x0
    for header, w in zip(headers, widths):
        add_panel(slide, x, y_head, w, 0.52, fill=ACCENT_FILL, line=BORDER)
        add_textbox(slide, x + 0.08, y_head + 0.13, w - 0.16, 0.22, header, 13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        x += w + 0.08

    rows = [
        ("2", "0.174 m²", "0.087 m²", "0.87 kg", "44.7 km", "efficient"),
        ("4", "0.521 m²", "0.260 m²", "2.60 kg", "103.9 km", "best balance"),
        ("8", "1.215 m²", "0.607 m²", "6.07 kg", "162.7 km", "aggressive"),
    ]
    fills = [RGBColor(254, 249, 195), RGBColor(255, 247, 237), RGBColor(254, 242, 242)]
    for ridx, row in enumerate(rows):
        y = y_head + 0.68 + ridx * row_h
        x = x0
        for cidx, (value, w) in enumerate(zip(row, widths)):
            add_panel(slide, x, y, w, 0.58, fill=fills[ridx], line=BORDER)
            add_textbox(
                slide,
                x + 0.08,
                y + 0.15,
                w - 0.16,
                0.22,
                value,
                14,
                bold=(cidx == 0 or (ridx == 1 and cidx in (3, 4, 5))),
                color=TITLE if ridx == 1 else TEXT,
                align=PP_ALIGN.CENTER,
            )
            x += w + 0.08

    add_panel(slide, 0.78, 5.88, 5.55, 0.9)
    add_textbox(slide, 0.98, 6.04, 5.1, 0.2, "Key read", 15, bold=True, color=ACCENT)
    add_bullets(
        slide,
        0.98,
        6.28,
        5.05,
        0.28,
        [
            "β ratio ≈ 4 is the current presentation case because it is strong without demanding > 6 kg of panel-system mass.",
        ],
        size=13,
    )
    add_panel(slide, 6.65, 5.88, 5.75, 0.9)
    add_textbox(slide, 6.85, 6.04, 5.3, 0.2, "Mass assumption", 15, bold=True, color=ACCENT)
    add_textbox(
        slide,
        6.85,
        6.28,
        5.25,
        0.3,
        "* Added panel-system mass assumes 5 kg/m² areal density for panels + hinges + deployment hardware.",
        12,
        color=TEXT,
    )


def slide9(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(
        slide,
        "Backup: Open Risks And Next Gating Questions",
        "The main remaining questions are engineering closure questions, not more first-cut trajectory plots.",
    )
    add_panel(slide, 0.58, 1.88, 5.95, 4.98)
    add_textbox(slide, 0.85, 2.1, 5.35, 0.3, "What still needs to close", 20, bold=True, color=ACCENT)
    add_bullets(
        slide,
        0.84,
        2.46,
        5.3,
        4.1,
        [
            "Mass closure: panel area helps only if the full panel + hinge + release + local structure mass stays favorable.",
            "Trim and stability: deployment changes moments and center-of-pressure location, not just drag.",
            "Thermal and structural survivability: panel heating, attachment loads, and clean jettison still need to be shown.",
            "Impact-system compatibility: added surfaces must not break the hard-lander impact concept.",
        ],
        size=15,
    )
    add_panel(slide, 6.78, 1.88, 5.95, 4.98)
    add_textbox(slide, 7.05, 2.1, 5.35, 0.3, "What the next iteration should answer", 20, bold=True, color=ACCENT)
    add_bullets(
        slide,
        7.04,
        2.46,
        5.3,
        4.1,
        [
            "How much authority do we gain per added kilogram and per added square meter once hardware packaging is included?",
            "Can the deployed configuration be kept trim-stable enough to preserve the low-cost control story?",
            "Do winds and navigation uncertainty materially erode the demonstrated downrange benefit?",
            "Would more than 2 panels improve crossrange without destroying the downrange drag-control benefit?",
        ],
        size=15,
    )
    add_panel(slide, 0.72, 6.15, 12.0, 0.62, fill=ACCENT_FILL, line=BORDER)
    add_textbox(
        slide,
        0.96,
        6.34,
        11.5,
        0.18,
        "Current proposal-safe claim: strong preliminary downrange authority is demonstrated; crossrange is promising but still architecture-limited.",
        14,
        bold=True,
        color=ACCENT,
    )
    set_slide_notes(
        slide,
        "This backup slide is the clean way to close the technical story. "
        "The main point is that the open issues are now engineering closure questions: mass, trim, structure, heating, and packaging. "
        "That is a good place to be for a proposal, because it means the first-cut authority result is already strong enough to motivate the next step.",
    )


def build_outline() -> None:
    outline = """# AERO-08 Tuesday Slide Outline

## Slide 1. Low-Cost Guided Mars Surface Access
- Figure: no heavy plot; concept framing plus representative authority callout
- Takeaways:
  - SHIELD-like hard-landers are low-cost, but accept large landing footprints.
  - The concept is discrete aerodynamic area modulation during entry.
  - The goal is coarse targeting and footprint reduction, not precision landing.
- Speaker note:
  - Open with the mission gap and the proposal framing.

## Slide 2. Drag Modulation Produces Meaningful Downrange Authority
- Figure: `downrange_authority_heatmap.png`
- Takeaways:
  - Downrange authority is real.
  - Authority grows with `β_ratio`.
  - Early jettison collapses the benefit.
- Speaker note:
  - This is the main answer to the “5 km or 50 km?” question.

## Slide 3. Authority Is Generated In A Specific Entry Corridor
- Figure: `local_control_effectiveness.png`
- Support figure: `drag_accel_vs_altitude.png`
- Takeaways:
  - Control is concentrated in a mid-entry corridor.
  - The effect is physically interpretable, not a numerical artifact.
  - Discrete switch-altitude logic is a natural guidance choice.
- Speaker note:
  - Explain where the authority comes from. In the support plot, 'Body only' means no deployed surfaces, 'Always deployed until 20 km' means continuous deployment until jettison, and 'Deploy below X km, jettison at 20 km' means the surfaces are activated only below the named switch altitude.

## Slide 4. Control Authority Comes With A Design Trade
- Figure: `terminal_velocity_and_peak_g_vs_beta_ratio.png`
- Takeaways:
  - Slides 2-4 are symmetric drag-modulation results, so they quantify downrange capability.
  - More area buys authority, but changes terminal conditions and loads.
  - `β_ratio ≈ 4` is the current sweet spot.
  - Jettison above roughly `20–30 km` hurts badly.
- Speaker note:
  - Use this slide to defend the disciplined minimum-viable architecture and to state clearly that the main demonstrated result is downrange control from drag modulation.

## Slide 5. Crossrange Requires Asymmetric Deployment
- Figure: `crossrange_sensitivity.png`
- Inset: `landing_ellipse_centered.png`
- Takeaways:
  - Crossrange is not part of the symmetric drag-only results on slides 2-4.
  - Asymmetric deployment can induce crossrange.
  - With 2 panels, crossrange is still strongly coupled to downrange.
  - Additional panels may materially change the control allocation.
- Speaker note:
  - Close without overselling 3D guidance.

## Slide 6. Backup: Modeling Setup And Initial Conditions
- Purpose:
  - State clearly what the current model includes and excludes.
- Main points:
  - MarsGRAM-backed 3DOF point-mass model
  - Nominal entry state: `h0 = 125 km`, `V0 = 5.5 km/s`, `γ0 = -12°`
  - Main deterministic authority plots use zero winds; the 3D Monte Carlo ellipse uses MarsGRAM winds
  - Not yet a 6DOF or transition-flow truth model

## Slide 7. Backup: Vehicle And Aerodynamic Construction
- Purpose:
  - Show the SHIELD-informed surrogate body and deployable panels.
- Main points:
  - 70° sphere-cone surrogate
  - 2 symmetric deployable flat plates
  - Modified Newtonian body, broadside drag panels

## Slide 8. Backup: What β Ratio Means In Hardware
- Purpose:
  - Translate β ratio into actual panel size and added panel-system mass.
- Main points:
  - β ratio = β_high / β_low
  - For the nominal 150 kg/m² case:
    - β ratio 2 -> 0.174 m² total panel area, 0.87 kg added mass
    - β ratio 4 -> 0.521 m² total panel area, 2.60 kg added mass
    - β ratio 8 -> 1.215 m² total panel area, 6.07 kg added mass

## Slide 9. Backup: Open Risks And Next Gating Questions
- Purpose:
  - Show that the next uncertainties are engineering closure questions, not a lack of preliminary authority.
- Main points:
  - Mass, trim, thermal, structural, and jettison closure remain open
  - Winds are already separated correctly from the main deterministic authority claim
  - The next useful step is to connect the demonstrated authority to hardware closure and richer crossrange architectures
"""
    OUTLINE_PATH.write_text(outline)


def export_quicklook_preview() -> None:
    subprocess.run(
        ["qlmanage", "-p", "-o", str(QLPREVIEW_DIR.parent), str(PPTX_PATH)],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    generated = QLPREVIEW_DIR.parent / f"{PPTX_PATH.name}.qlpreview"
    if generated.exists() and generated != QLPREVIEW_DIR:
        if QLPREVIEW_DIR.exists():
            shutil.rmtree(QLPREVIEW_DIR)
        generated.rename(QLPREVIEW_DIR)


def main() -> None:
    clean_outputs()
    prs = Presentation()
    prs.slide_width = inches(13.333)
    prs.slide_height = inches(7.5)
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    slide9(prs)
    prs.save(str(PPTX_PATH))
    build_outline()
    export_quicklook_preview()
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
